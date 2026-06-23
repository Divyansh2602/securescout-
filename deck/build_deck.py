#!/usr/bin/env python3
"""Builds the SecureScout pitch deck (16:9) as a .pptx.

Run:  python deck/build_deck.py
Out:  deck/SecureScout_PSB_Hackathon_2026.pptx

Pure python-pptx, no template needed. Re-run any time to regenerate.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- brand palette ---------------------------------------------------------
BG      = RGBColor(0x0B, 0x0D, 0x14)
PANEL   = RGBColor(0x14, 0x18, 0x28)
INDIGO  = RGBColor(0x81, 0x8C, 0xF8)
TEXT    = RGBColor(0xC8, 0xD0, 0xE7)
MUTED   = RGBColor(0x8A, 0x93, 0xAD)
WHITE   = RGBColor(0xF0, 0xF3, 0xFF)
CRIT    = RGBColor(0xEF, 0x44, 0x44)
HIGH    = RGBColor(0xF9, 0x73, 0x16)
MED     = RGBColor(0xEA, 0xB3, 0x08)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)  # send to back
    return s


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0):
    sp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (txt,size,color,bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            run = p.add_run(); run.text = txt
            run.font.size = Pt(size); run.font.color.rgb = color
            run.font.bold = bold; run.font.name = "Segoe UI"
    return tb


def kicker(s, label):
    text(s, 0.7, 0.55, 11, 0.4, [[(label.upper(), 12, INDIGO, True)]])


def title(s, t):
    text(s, 0.7, 0.9, 12, 1.0, [[(t, 32, WHITE, True)]])


# ===========================================================================
# 1 — Title
# ===========================================================================
s = slide()
box(s, 0, 0, 13.333, 0.18, fill=INDIGO)
text(s, 0.9, 2.3, 11.5, 1.2, [[("🛡  SecureScout", 54, WHITE, True)]])
text(s, 0.95, 3.5, 11.5, 0.8,
     [[("Automated discovery of misconfigurations & vulnerable dependencies", 20, TEXT, False)]])
text(s, 0.95, 4.2, 11.5, 0.6,
     [[("in open-source software — with one-click auto-remediation.", 20, TEXT, False)]])
box(s, 0.95, 5.3, 5.6, 1.3, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
text(s, 1.2, 5.45, 5.2, 1.1, [
    [("PSB Hackathon 2026", 16, INDIGO, True)],
    [("Problem Statement 3 · UCO Bank × IIT Kharagpur", 12.5, MUTED, False)],
    [("Divyansh Gupta · B.Tech CSE, VIT", 12.5, MUTED, False)],
], space_after=4)
text(s, 7.0, 5.55, 5.4, 1.1, [
    [("49", 30, CRIT, True), ("  findings on the demo target", 13, MUTED, False)],
    [("100/100", 24, HIGH, True), ("  risk score · ", 13, MUTED, False),
     ("14", 18, GREEN, True), ("  auto-patched", 13, MUTED, False)],
], space_after=6)

# ===========================================================================
# 2 — Problem
# ===========================================================================
s = slide(); kicker(s, "The Problem"); title(s, "Vulnerabilities hide in what you didn't write")
pts = [
    ("80%+ of modern app code", "comes from open-source dependencies — each a potential entry point."),
    ("Misconfigurations & secrets", "hardcoded keys, weak crypto, debug mode, and exposed credentials ship to production."),
    ("Known CVEs go unpatched", "teams lack a fast, automated way to find AND fix vulnerable package versions."),
    ("Banking-grade risk", "a single unpatched dependency can expose customer financial data."),
]
y = 2.1
for head, body in pts:
    box(s, 0.7, y, 0.12, 0.95, fill=CRIT)
    text(s, 1.0, y, 11.4, 1.0, [
        [(head, 18, WHITE, True)],
        [(body, 14, MUTED, False)],
    ], space_after=2)
    y += 1.15

# ===========================================================================
# 3 — Solution
# ===========================================================================
s = slide(); kicker(s, "The Solution"); title(s, "Scan → Score → Report → Auto-Fix, in seconds")
cards = [
    ("Scan", "5 detectors sweep deps, code, configs & secrets across the whole project."),
    ("Score", "A single 0–100 risk score with PASS/FAIL — built for CI gates."),
    ("Report", "Premium self-contained HTML + JSON, plus a live web dashboard."),
    ("Auto-Fix", "Generates a patched requirements.fixed.txt with safe version pins."),
]
x = 0.7
for i, (h, b) in enumerate(cards):
    box(s, x, 2.2, 2.85, 3.0, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
    text(s, x+0.25, 2.45, 2.4, 0.5, [[(f"0{i+1}", 16, INDIGO, True)]])
    text(s, x+0.25, 2.95, 2.4, 0.6, [[(h, 20, WHITE, True)]])
    text(s, x+0.25, 3.6, 2.4, 1.5, [[(b, 13, MUTED, False)]], line_spacing=1.05)
    x += 3.05

# ===========================================================================
# 4 — Architecture
# ===========================================================================
s = slide(); kicker(s, "Architecture"); title(s, "One engine, five detectors, two surfaces")
# detector row
dets = [("CVE", "Vulnerable\ndependencies"), ("Secret", "Hardcoded\nkeys/passwords"),
        ("Crypto", "Weak hashing\n& ciphers"), ("Injection", "SQLi / cmd /\nXSS / eval"),
        ("Config", "Insecure\nsettings")]
x = 0.7
for name, desc in dets:
    box(s, x, 2.15, 2.35, 1.3, fill=PANEL, line=INDIGO, line_w=1.0)
    text(s, x, 2.3, 2.35, 0.5, [[(name, 16, INDIGO, True)]], align=PP_ALIGN.CENTER)
    text(s, x, 2.78, 2.35, 0.6,
         [[(d, 11, MUTED, False)] for d in [desc.replace("\n", " ")]], align=PP_ALIGN.CENTER)
    x += 2.5
text(s, 0.7, 3.7, 12, 0.4, [[("⮧  unified Scanner Engine  (dedupe · severity rank · risk score)", 14, TEXT, True)]],
     align=PP_ALIGN.CENTER)
# two surfaces
box(s, 1.7, 4.5, 4.4, 1.7, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
text(s, 1.95, 4.65, 4.0, 1.5, [
    [("CLI  +  CI/CD", 17, WHITE, True)],
    [("Python CLI · JSON/HTML export", 12.5, MUTED, False)],
    [("--autofix · --fail-on-critical", 12.5, MUTED, False)],
    [("GitHub Actions gate", 12.5, MUTED, False)],
], space_after=3)
box(s, 7.2, 4.5, 4.4, 1.7, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
text(s, 7.45, 4.65, 4.0, 1.5, [
    [("Enterprise Platform", 17, WHITE, True)],
    [("Next.js 14 dashboard + upload", 12.5, MUTED, False)],
    [("Express API · JWT · 5-tier RBAC", 12.5, MUTED, False)],
    [("PostgreSQL + Prisma · audit log", 12.5, MUTED, False)],
], space_after=3)

# ===========================================================================
# 5 — Live results
# ===========================================================================
s = slide(); kicker(s, "Live Demo"); title(s, "vulnerable_bank_app — caught in 0.01s")
stats = [("49", "Total findings", TEXT), ("20", "Critical", CRIT),
         ("23", "High", HIGH), ("6", "Medium", MED)]
x = 0.7
for num, lbl, c in stats:
    box(s, x, 2.2, 2.85, 1.7, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
    text(s, x, 2.45, 2.85, 0.9, [[(num, 40, c, True)]], align=PP_ALIGN.CENTER)
    text(s, x, 3.35, 2.85, 0.4, [[(lbl, 13, MUTED, False)]], align=PP_ALIGN.CENTER)
    x += 3.05
box(s, 0.7, 4.25, 11.93, 1.9, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
text(s, 1.0, 4.45, 11.3, 1.7, [
    [("Risk Score 100/100 · Status FAIL", 18, CRIT, True)],
    [("Catches CVEs (PyYAML RCE, Paramiko auth bypass), hardcoded AWS keys & JWT secrets,", 14, TEXT, False)],
    [("SQL injection, command injection, eval/exec, MD5/SHA-1, pickle, debug mode, ECB crypto…", 14, TEXT, False)],
    [("→ Auto-Fix patched 14 vulnerable dependencies to safe versions automatically.", 14, GREEN, True)],
], space_after=6)

# ===========================================================================
# 6 — Auto-Fix
# ===========================================================================
s = slide(); kicker(s, "Differentiator"); title(s, "It doesn't just find — it fixes")
text(s, 0.7, 2.0, 12, 0.5,
     [[("--autofix generates a patched requirements file with known-safe pins:", 15, TEXT, False)]])
box(s, 0.7, 2.7, 11.93, 2.6, fill=RGBColor(0x0E,0x12,0x1F), line=RGBColor(0x1E,0x22,0x36))
mono = [
    ("flask==0.12.2", "flask>=2.3.3", "CVE-2018-1000656"),
    ("pyyaml==3.12", "pyyaml>=6.0.1", "CVE-2017-18342"),
    ("paramiko==2.4.0", "paramiko>=3.3.1", "CVE-2018-1000805"),
    ("cryptography==1.9", "cryptography>=41.0.5", "CVE-2018-10903"),
    ("…+10 more", "", ""),
]
y = 2.9
for old, new, cve in mono:
    para = [[(f"{old:<22}", 13, CRIT, False)]]
    line = [(f"{old:<22}", 13, CRIT, False)]
    if new:
        line += [("→  ", 13, MUTED, False), (f"{new:<22}", 13, GREEN, True), (f"  {cve}", 12, MUTED, False)]
    text(s, 1.0, y, 11.3, 0.4, [line])
    y += 0.42
text(s, 0.7, 5.6, 12, 0.8,
     [[("Code-level issues (secrets, injection) are flagged for manual review — never auto-patched,", 13, MUTED, False)],
      [("so application logic is never silently broken.", 13, MUTED, False)]], space_after=2)

# ===========================================================================
# 7 — Enterprise platform
# ===========================================================================
s = slide(); kicker(s, "Production-Grade"); title(s, "A full enterprise platform, not a script")
rows = [
    ("Frontend", "Next.js 14 · TypeScript · Tailwind · shadcn/ui · drag-drop ZIP upload · live dashboard"),
    ("API", "Express + TypeScript · JWT w/ refresh tokens · 5-tier RBAC · Zod validation · Helmet"),
    ("Security", "3-tier rate limiting · audit logging · unit + integration tests"),
    ("Data", "PostgreSQL + Prisma ORM"),
    ("Deploy (no Docker)", "Vercel · Render · Neon/Supabase · Upstash — GitHub Actions CI/CD"),
]
y = 2.1
for k, v in rows:
    box(s, 0.7, y, 3.0, 0.78, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
    text(s, 0.85, y+0.14, 2.8, 0.5, [[(k, 14, INDIGO, True)]])
    text(s, 3.95, y+0.1, 8.7, 0.6, [[(v, 13, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.92

# ===========================================================================
# 8 — CI/CD
# ===========================================================================
s = slide(); kicker(s, "Shift Left"); title(s, "Blocks vulnerable code at the PR")
box(s, 0.7, 2.2, 11.93, 1.4, fill=RGBColor(0x0E,0x12,0x1F), line=RGBColor(0x1E,0x22,0x36))
text(s, 1.0, 2.4, 11.3, 1.1, [
    [("$ securescout ./app --fail-on-critical", 16, GREEN, True)],
    [("Risk Score: 100/100 · Status: FAIL", 14, CRIT, False)],
    [("Error: 20 critical findings — build failed ✗", 14, TEXT, False)],
], space_after=4)
text(s, 0.7, 4.0, 12, 2.0, [
    [("• Runs on every push & pull request via GitHub Actions", 15, TEXT, False)],
    [("• Non-zero exit on critical findings → red build, merge blocked", 15, TEXT, False)],
    [("• Uploads JSON + HTML report as build artifacts", 15, TEXT, False)],
    [("• Security becomes a gate, not an afterthought", 15, INDIGO, True)],
], space_after=10)

# ===========================================================================
# 9 — Why it wins
# ===========================================================================
s = slide(); kicker(s, "Impact"); title(s, "Why SecureScout stands out")
wins = [
    ("End-to-end", "Detection AND remediation — most tools stop at detection."),
    ("Demo-ready", "One command → 49 real findings, premium report, auto-fix."),
    ("Enterprise depth", "Real auth, RBAC, audit, CI/CD — production architecture."),
    ("Domain fit", "Built for banking-grade dependency risk (PS3)."),
]
x, y = 0.7, 2.2
for i, (h, b) in enumerate(wins):
    px = 0.7 + (i % 2) * 6.1
    py = 2.2 + (i // 2) * 2.0
    box(s, px, py, 5.85, 1.7, fill=PANEL, line=RGBColor(0x1E,0x22,0x36))
    text(s, px+0.3, py+0.2, 5.3, 0.5, [[(h, 19, INDIGO, True)]])
    text(s, px+0.3, py+0.8, 5.3, 0.8, [[(b, 14, TEXT, False)]], line_spacing=1.05)

# ===========================================================================
# 10 — Thank you
# ===========================================================================
s = slide()
box(s, 0, 0, 13.333, 0.18, fill=INDIGO)
text(s, 0.9, 2.7, 11.5, 1.2, [[("Thank you.", 46, WHITE, True)]])
text(s, 0.95, 3.9, 11.5, 0.6, [[("SecureScout — find it, score it, fix it.", 20, TEXT, False)]])
text(s, 0.95, 4.9, 11.5, 1.0, [
    [("Divyansh Gupta · divyansh.gupta2023@vitstudent.ac.in", 14, MUTED, False)],
    [("PSB Hackathon 2026 · Problem Statement 3", 14, MUTED, False)],
], space_after=4)

# ---- save ------------------------------------------------------------------
out = Path(__file__).parent / "SecureScout_PSB_Hackathon_2026.pptx"
prs.save(str(out))
print(f"Deck written: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
